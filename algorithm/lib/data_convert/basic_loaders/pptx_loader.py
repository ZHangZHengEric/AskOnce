import os
import hashlib
from pptx import Presentation
from AskOnce.algorithm.lib.data_convert.basic_loaders.base_loader import BaseLoader

class PptxLoader(BaseLoader):
    def load_data(self, url):
        # meta_data = {"url": url, "file_size": f"{os.path.getsize(url)/ 1024 **2:.02f}MB", "file_type": url.split(".")[-1]}
        pr = Presentation(url)
        data = []
        text = ""
        for slide in pr.slides:
            content = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content += shape.text
            text += content
            data.append({"content":content})
        doc_id = hashlib.sha256((text + url).encode()).hexdigest()
        text_details = {"doc_id":doc_id, "data": data}
        return text, text_details